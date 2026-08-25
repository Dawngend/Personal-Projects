"""One-command staging load and OCR/grading smoke suite."""

from __future__ import annotations

from concurrent.futures import ThreadPoolExecutor
import json
import os
from pathlib import Path
import re
import statistics
import tempfile
import time
from urllib.request import Request, urlopen

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches
import pytesseract

from extractor import process_module_file_v2
from grading import grade_problem_answer


CASES = {
    "CASE18": ("{lambda = 5, 1}", "{1, 5}", "structured"),
    "CASE20": ("t(1, -1, 0)", "(t, -t, 0)", "symbolic"),
    "CASE21": ("s(2, 4) + t(1, 1)", "(2s+t, 4s+t)", "symbolic"),
    "CASE30": ("x = 1, y = 2", "(1, 2)", "structured"),
}


def request_json(url: str) -> dict[str, object]:
    with urlopen(Request(url, headers={"User-Agent": "andyhub-staging-smoke"}), timeout=20) as response:
        if response.status != 200:
            raise RuntimeError(f"{url} returned HTTP {response.status}")
        return json.loads(response.read().decode("utf-8"))


def load_smoke(base_url: str, requests: int = 24, concurrency: int = 6) -> dict[str, object]:
    for path, service in (
        ("/proxy-health", "proxy"),
        ("/health", "web"),
        ("/api/v1/health", "api"),
    ):
        payload = request_json(base_url + path)
        if payload.get("status") != "ok" or payload.get("service") != service:
            raise RuntimeError(f"unexpected {service} health payload: {payload}")
    worker_health = request_json(os.environ.get("WORKER_HEALTH_URL", "http://worker:8001/health"))
    if worker_health.get("status") != "ok":
        raise RuntimeError(f"unexpected worker health payload: {worker_health}")

    def load_once(_: int) -> float:
        started = time.perf_counter()
        with urlopen(base_url + "/", timeout=20) as response:
            body = response.read().decode("utf-8", errors="replace")
            if response.status != 200 or "AndyHub" not in body:
                raise RuntimeError("workspace load did not return the AndyHub page")
        return time.perf_counter() - started

    with ThreadPoolExecutor(max_workers=concurrency) as pool:
        durations = list(pool.map(load_once, range(requests)))
    ordered = sorted(durations)
    p95 = ordered[max(0, int(len(ordered) * 0.95) - 1)]
    maximum = float(os.environ.get("LOAD_P95_MAX_SECONDS", "3"))
    if p95 > maximum:
        raise RuntimeError(f"workspace p95 {p95:.3f}s exceeds {maximum:.3f}s")
    return {
        "requests": requests,
        "concurrency": concurrency,
        "p50_seconds": round(statistics.median(durations), 3),
        "p95_seconds": round(p95, 3),
    }


def ocr_grading_smoke() -> dict[str, object]:
    pytesseract.get_tesseract_version()
    with tempfile.TemporaryDirectory() as temporary_directory:
        root = Path(temporary_directory)
        image_path = root / "linear-algebra-smoke.png"
        module_path = root / "linear-algebra-smoke.pptx"
        image = Image.new("RGB", (2200, 1300), "white")
        draw = ImageDraw.Draw(image)
        font = ImageFont.truetype(
            os.environ.get("SMOKE_FONT", "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
            58,
        )
        lines = [f"{case_id}: {values[0]}" for case_id, values in CASES.items()]
        draw.multiline_text((90, 100), "\n\n".join(lines), fill="black", font=font, spacing=36)
        image.save(image_path)

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(image_path), Inches(0.2), Inches(0.2), width=Inches(12.8))
        presentation.save(module_path)

        extracted = process_module_file_v2(str(module_path))
        results: dict[str, str] = {}
        for case_id, (expected_input, expected_answer, expected_tier) in CASES.items():
            match = re.search(rf"(?m)^{case_id}:\s*(.+?)\s*$", extracted)
            if not match:
                raise RuntimeError(f"OCR output did not contain {case_id}: {extracted!r}")
            ocr_input = match.group(1)
            if ocr_input != expected_input:
                raise RuntimeError(
                    f"OCR changed {case_id}: expected {expected_input!r}, got {ocr_input!r}"
                )
            result = grade_problem_answer(ocr_input, expected_answer)
            if not result or result.tier != expected_tier:
                raise RuntimeError(
                    f"grading failed for {case_id}: input={ocr_input!r}, "
                    f"expected={expected_answer!r}, tier={result.tier!r}"
                )
            results[case_id] = result.tier
    return results


def main() -> int:
    base_url = os.environ.get("SMOKE_BASE_URL", "http://proxy:8080").rstrip("/")
    result = {
        "status": "ok",
        "load": load_smoke(base_url),
        "ocr_grading": ocr_grading_smoke(),
    }
    print(json.dumps(result, indent=2, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
